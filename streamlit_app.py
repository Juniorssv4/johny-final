import streamlit as st
import time
import google.generativeai as genai
import requests
from io import BytesIO
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from tenacity import retry, stop_after_attempt, wait_exponential, RetryError
import json
import base64

# ───────────────────────────────────────────────
# SECRETS & CONFIG
# ───────────────────────────────────────────────
try:
    genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
except:
    st.error("Add GEMINI_API_KEY in Secrets")
    st.stop()

GITHUB_TOKEN = st.secrets["GITHUB_TOKEN"]
GITHUB_REPO = st.secrets["GITHUB_REPO"]
GLOSSARY_FILE = "glossary.txt"

# Primary and fallback models
PRIMARY_MODEL = "gemini-2.5-flash"
FALLBACK_MODEL = "gemini-1.5-flash"

if "current_model" not in st.session_state:
    st.session_state.current_model = PRIMARY_MODEL

model = genai.GenerativeModel(st.session_state.current_model)

# Backoff for rate limits
@retry(
    stop=stop_after_attempt(6),
    wait=wait_exponential(multiplier=1, min=4, max=60)
)
def safe_generate_content(prompt):
    return model.generate_content(prompt)

# ───────────────────────────────────────────────
# GLOSSARY: LOAD FROM GITHUB + EDIT IN-APP + SAVE BACK TO GITHUB
# ───────────────────────────────────────────────
def load_glossary():
    try:
        url = f"https://raw.githubusercontent.com/{GITHUB_REPO}/main/{GLOSSARY_FILE}"
        response = requests.get(url)
        response.raise_for_status()
        lines = response.text.splitlines()
        glossary_dict = {}
        for line in lines:
            line = line.strip()
            if line and ":" in line:
                parts = line.split(":", 1)
                eng = parts[0].strip().lower()
                lao = parts[1].strip() if len(parts) > 1 else ""
                glossary_dict[eng] = lao
        return glossary_dict
    except Exception as e:
        st.error(f"Failed to load glossary: {str(e)}")
        return {}

def save_glossary_to_github(glossary_dict):
    try:
        # Get current file SHA
        url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{GLOSSARY_FILE}"
        headers = {"Authorization": f"token {GITHUB_TOKEN}"}
        response = requests.get(url, headers=headers)
        sha = response.json().get('sha') if response.status_code == 200 else None

        # Build file content
        lines = [f"{eng.capitalize()}: {lao}" for eng, lao in sorted(glossary_dict.items())]
        content = "\n".join(lines) + "\n"
        encoded_content = base64.b64encode(content.encode()).decode()

        # Commit
        data = {
            "message": "Update glossary from app",
            "content": encoded_content,
            "sha": sha
        }
        requests.put(url, headers=headers, json=data)
        st.success("Glossary saved to GitHub!")
        st.session_state.glossary = glossary_dict  # update memory
        st.rerun()
    except Exception as e:
        st.error(f"Failed to save glossary: {str(e)}")

# Load glossary on every run (no cache)
glossary = load_glossary()

def get_glossary_prompt():
    if glossary:
        terms = "\n".join([f"• {e.capitalize()} → {l}" for e, l in glossary.items()])
        return f"Use EXACTLY these terms:\n{terms}\n"
    return ""

def translate_text(text, direction):
    if not text.strip():
        return ""
    target = "Lao" if direction == "English → Lao" else "English"
    prompt = f"""{get_glossary_prompt()}Translate ONLY the text to {target}.
Return ONLY the translation.
Text: {text}"""
    try:
        response = safe_generate_content(prompt)
        return response.text.strip()
    except RetryError as e:
        if "429" in str(e.last_attempt.exception()) or "quota" in str(e.last_attempt.exception()).lower():
            if st.session_state.current_model == PRIMARY_MODEL:
                st.session_state.current_model = FALLBACK_MODEL
                st.info("Rate limit on gemini-2.5-flash — switched to gemini-1.5-flash.")
                global model
                model = genai.GenerativeModel(FALLBACK_MODEL)
                response = model.generate_content(prompt)
                return response.text.strip()
        st.error("Timed out after retries — try again in 5 minutes.")
        return "[Failed — try later]"
    except Exception as e:
        st.error(f"API error: {str(e)}")
        return "[Failed — try again]"

# UI
st.set_page_config(
    page_title="Johny",
    page_icon="https://raw.githubusercontent.com/Juniorssv4/johny-final/main/Johny.png",
    layout="centered"
)

st.title("😊 Johny — NPA Lao Translator")

direction = st.radio("Direction", ["English → Lao", "Lao → English"], horizontal=True)

tab1, tab2 = st.tabs(["Translate Text", "Translate File"])

with tab1:
    text = st.text_area("Enter text to translate", height=200)
    if st.button("Translate Text", type="primary"):
        with st.spinner("Translating..."):
            result = translate_text(text, direction)
            st.success("Translation:")
            st.write(result)

with tab2:
    uploaded_file = st.file_uploader("Upload DOCX • XLSX • PPTX (max 50MB)", type=["docx", "xlsx", "pptx"])
    if uploaded_file:
        MAX_SIZE_MB = 50
        if uploaded_file.size > MAX_SIZE_MB * 1024 * 1024:
            st.error(f"File too large! Max allowed size is {MAX_SIZE_MB}MB. Your file is {uploaded_file.size / (1024*1024):.1f}MB.")
        elif st.button("Translate File", type="primary"):
            with st.spinner("Translating file..."):
                file_bytes = uploaded_file.read()
                file_name = uploaded_file.name
                ext = file_name.rsplit(".", 1)[-1].lower()
                output = BytesIO()
                total_elements = 0
                elements_list = []
                if ext == "docx":
                    doc = Document(BytesIO(file_bytes))
                    for p in doc.paragraphs:
                        if p.text.strip():
                            total_elements += 1
                            elements_list.append(("para", p))
                    for table in doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                for p in cell.paragraphs:
                                    if p.text.strip():
                                        total_elements += 1
                                        elements_list.append(("para", p))
                elif ext == "xlsx":
                    wb = load_workbook(BytesIO(file_bytes))
                    for ws in wb.worksheets:
                        for row in ws.iter_rows():
                            for cell in row:
                                if isinstance(cell.value, str) and cell.value.strip():
                                    total_elements += 1
                                    elements_list.append(("cell", cell))
                elif ext == "pptx":
                    prs = Presentation(BytesIO(file_bytes))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for p in shape.text_frame.paragraphs:
                                    if p.text.strip():
                                        total_elements += 1
                                        elements_list.append(("para", p))
                if total_elements == 0:
                    st.warning("No text found in file.")
                    st.stop()
                progress_bar = st.progress(0)
                status_text = st.empty()
                translated_count = 0
                for element_type, element in elements_list:
                    status_text.text(f"Translating... {translated_count}/{total_elements}")
                    if element_type == "para":
                        translated = translate_text(element.text, direction)
                        element.text = translated
                    elif element_type == "cell":
                        translated = translate_text(element.value, direction)
                        element.value = translated
                    translated_count += 1
                    progress_bar.progress(translated_count / total_elements)
                status_text.text("Saving file...")
                if ext == "docx":
                    doc.save(output)
                elif ext == "xlsx":
                    wb.save(output)
                elif ext == "pptx":
                    prs.save(output)
                output.seek(0)
                filename = f"TRANSLATED_{file_name}"
                mime_type = "application/octet-stream"
                st.success("Translation complete!")
                st.info("Click the big button below to download your translated file. Your browser may block auto-downloads — this button always works!")
                st.download_button(
                    label="📥 DOWNLOAD TRANSLATED FILE NOW",
                    data=output,
                    file_name=filename,
                    mime=mime_type,
                    type="primary",
                    use_container_width=True,
                    key="download_btn_" + str(time.time()),
                    help="Click here to save the translated file to your device"
                )
                st.caption("Tip: If nothing happens, refresh the page or try in another browser (Chrome works best).")

# Teach term – EDIT IN-APP + SAVE TO GITHUB
with st.expander("➕ Teach Johny a new term (edit in app & save to GitHub)"):
    st.info("Add or edit terms directly here. Changes are saved to glossary.txt in GitHub.")

    # Show current glossary
    if glossary:
        st.write("Current terms:")
        st.json(glossary)
    else:
        st.info("No terms loaded yet.")

    # Form to add new term
    new_eng = st.text_input("English term")
    new_lao = st.text_input("Lao translation")
    if st.button("Add / Update Term"):
        if new_eng and new_lao:
            glossary[new_eng.strip().lower()] = new_lao.strip()
            # Save to GitHub
            lines = [f"{eng.capitalize()}: {lao}" for eng, lao in sorted(glossary.items())]
            content = "\n".join(lines) + "\n"
            url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{GLOSSARY_FILE}"
            headers = {"Authorization": f"token {GITHUB_TOKEN}", "Content-Type": "application/json"}
            response = requests.get(url, headers=headers)
            sha = response.json().get('sha') if response.status_code == 200 else None
            encoded = base64.b64encode(content.encode()).decode()
            data = {
                "message": "Update glossary from app",
                "content": encoded,
                "sha": sha
            }
            put_response = requests.put(url, headers=headers, json=data)
            if put_response.status_code in (200, 201):
                st.success(f"Term '{new_eng}' → '{new_lao}' saved!")
                st.session_state.glossary = glossary
                st.rerun()
            else:
                st.error(f"GitHub save failed: {put_response.text}")
        else:
            st.error("Both fields required")

# Show count
st.caption(f"Active glossary: {len(glossary)} terms • Model: {st.session_state.current_model}")
